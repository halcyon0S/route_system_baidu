# merge_to_pdf.py
"""
将网组网点路线图中同一个工号的网组网点图合并到一张PDF中
命名规则：工号-姓名.pdf
保存到新文件夹：合并PDF

处理流程：
1. 检查"调整前"和"调整后"文件夹是否存在
2. 为每个存在的文件夹创建对应的封面页："工号-姓名-调整前"或"工号-姓名-调整后"
3. 将每个文件夹内的行政区图图片单独一页PDF
4. 将其余图片（网组网点图）2*2布局合并到PDF中
"""

import os
import re
from pathlib import Path
from typing import List, Tuple, Dict
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("⚠️ 警告：未安装 Pillow，将尝试使用 img2pdf")
    try:
        import img2pdf
        IMG2PDF_AVAILABLE = True
    except ImportError:
        IMG2PDF_AVAILABLE = False
        print("⚠️ 警告：未安装 img2pdf，请安装：pip install img2pdf")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ 警告：未安装 python-pptx，PPT功能将不可用")
    print("   如需使用PPT功能，请安装：pip install python-pptx")


def get_base_dir():
    """获取程序基础目录"""
    import sys
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    else:
        return os.path.dirname(os.path.abspath(__file__))


def parse_employee_folder(folder_name: str) -> Tuple[str, str] | None:
    """
    解析文件夹名称，提取工号和姓名
    格式：工号-姓名，例如：FJ10331281-陈成
    
    Returns:
        (employee_id, employee_name) 或 None
    """
    # 匹配格式：工号-姓名（工号和姓名之间用-分隔）
    match = re.match(r'^(.+?)-(.+)$', folder_name)
    if match:
        return match.group(1), match.group(2)
    return None


def is_district_map_file(filename: str) -> bool:
    """
    判断是否是行政区图文件
    格式：工号-姓名-行政区图_时间戳.png
    """
    return '行政区图' in filename


def is_group_map_file(filename: str) -> bool:
    """
    判断是否是网组网点图文件（排除行政区图）
    格式：工号-姓名-网组网点图-网组名_时间戳.png
    """
    # 包含"网组网点图"且不包含"行政区图"
    return '网组网点图' in filename and '行政区图' not in filename


def get_image_files_from_folder(folder_path: Path) -> Dict[str, List[Path]]:
    """
    从指定文件夹获取图片文件，分类为行政区图和网组网点图
    
    Args:
        folder_path: 文件夹路径（如：调整前）
        
    Returns:
        字典，包含 'district'（行政区图）和 'group'（网组网点图）两个列表
    """
    result = {
        'district': [],  # 行政区图
        'group': []      # 网组网点图
    }
    
    if not folder_path.exists() or not folder_path.is_dir():
        return result
    
    # 遍历文件夹中的所有PNG文件
    seen_files = set()  # 用于去重，避免重复添加同一文件
    for file in folder_path.iterdir():
        if file.is_file() and file.suffix.lower() == '.png':
            # 使用文件的绝对路径作为唯一标识，避免重复
            file_key = str(file.resolve())
            if file_key in seen_files:
                continue  # 跳过已处理的文件
            seen_files.add(file_key)
            
            if is_district_map_file(file.name):
                result['district'].append(file)
            elif is_group_map_file(file.name):
                result['group'].append(file)
    
    # 按文件名排序
    result['district'].sort(key=lambda x: x.name)
    result['group'].sort(key=lambda x: x.name)
    
    return result


def add_single_image_page(image_path: Path, pdf_pages: List[Image.Image]) -> bool:
    """
    将单张图片添加为一页PDF（用于行政区图）
    
    Args:
        image_path: 图片路径
        pdf_pages: PDF页面列表（会追加新页面）
        
    Returns:
        是否成功
    """
    if not PIL_AVAILABLE:
        return False
    
    try:
        # 打开图片
        try:
            img = Image.open(image_path)
        except Exception as open_error:
            print(f"  [错误] 无法打开图片 {image_path.name}: {open_error}")
            return False
        
        # 转换为RGB模式
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # PDF页面尺寸（A4横向 @ 300 DPI）
        PAGE_WIDTH = int(3508)
        PAGE_HEIGHT = int(2480)
        
        # 验证页面尺寸
        if PAGE_WIDTH <= 0 or PAGE_HEIGHT <= 0:
            print(f"  [错误] 无效的页面尺寸: {PAGE_WIDTH}x{PAGE_HEIGHT}")
            return False
        
        # 计算缩放比例，保持宽高比，适应页面
        img_width, img_height = img.size
        
        # 验证图片尺寸
        if img_width <= 0 or img_height <= 0:
            print(f"  [错误] 无效的图片尺寸: {img_width}x{img_height}, 文件: {image_path.name}")
            return False
        
        scale_w = PAGE_WIDTH / img_width
        scale_h = PAGE_HEIGHT / img_height
        scale = min(scale_w, scale_h)  # 使用较小的缩放比例，确保图片完全显示
        
        new_width = int(img_width * scale)
        new_height = int(img_height * scale)
        
        # 验证缩放后的尺寸
        if new_width <= 0 or new_height <= 0:
            print(f"  [错误] 缩放后尺寸无效: {new_width}x{new_height}, 原始: {img_width}x{img_height}")
            return False
        
        # 使用高质量缩放算法
        try:
            resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        except AttributeError:
            resized_img = img.resize((new_width, new_height), Image.LANCZOS)
        
        # 创建新页面（白色背景，确保尺寸是整数）
        try:
            page = Image.new('RGB', (PAGE_WIDTH, PAGE_HEIGHT), color='white')
        except Exception as e:
            print(f"  [错误] 创建页面失败: {e}")
            print(f"  [调试] PAGE_WIDTH: {PAGE_WIDTH} (类型: {type(PAGE_WIDTH)})")
            print(f"  [调试] PAGE_HEIGHT: {PAGE_HEIGHT} (类型: {type(PAGE_HEIGHT)})")
            import traceback
            traceback.print_exc()
            return False
        
        # 居中放置图片（确保位置是整数）
        center_x = int((PAGE_WIDTH - new_width) // 2)
        center_y = int((PAGE_HEIGHT - new_height) // 2)
        
        # 确保位置在页面范围内
        center_x = max(0, min(center_x, PAGE_WIDTH - new_width))
        center_y = max(0, min(center_y, PAGE_HEIGHT - new_height))
        
        page.paste(resized_img, (center_x, center_y))
        
        pdf_pages.append(page)
        print(f"  ✓ 添加单页图片: {image_path.name}")
        return True
        
    except Exception as e:
        print(f"  ⚠️ 添加单页图片失败 {image_path.name}: {e}")
        return False


def merge_images_to_pdf_pillow(image_paths: List[Path], output_path: Path) -> bool:
    """
    使用Pillow将图片合并为PDF，每页显示4张图片（2x2布局）
    
    Args:
        image_paths: 图片文件路径列表
        output_path: 输出PDF文件路径
        
    Returns:
        是否成功
    """
    if not PIL_AVAILABLE:
        return False
    
    try:
        images = []
        for img_path in image_paths:
            try:
                img = Image.open(img_path)
                # 转换为RGB模式（PDF需要RGB）
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                images.append(img)
                print(f"  ✓ 加载图片: {img_path.name}")
            except Exception as e:
                print(f"  ⚠️ 跳过无效图片 {img_path.name}: {e}")
                continue
        
        if not images:
            print("  ❌ 没有有效的图片可以合并")
            return False
        
        # 定义PDF页面尺寸（A4横向，适合每页2张上下布局）
        # 使用更大的尺寸以保持图片清晰度，不压缩
        # A4横向 @ 300 DPI 以获得更好的图片质量
        PAGE_WIDTH = int(3508)   # A4横向 @ 300 DPI (297mm * 300/25.4 ≈ 3508px)
        PAGE_HEIGHT = int(2480)  # A4横向 @ 300 DPI (210mm * 300/25.4 ≈ 2480px)
        IMAGES_PER_PAGE = 2
        GRID_COLS = 1
        GRID_ROWS = 2
        
        # 计算每张图片在页面中的尺寸（减小边距，特别是上下间隔不要太大）
        MARGIN_H = 10  # 左右边距
        MARGIN_V = 3   # 上下边距（进一步减小上下间隔，从5px改为3px）
        IMAGE_WIDTH = (PAGE_WIDTH - MARGIN_H * 2) // GRID_COLS  # 左右各10px
        IMAGE_HEIGHT = (PAGE_HEIGHT - MARGIN_V * 3) // GRID_ROWS  # 上下各3px，中间3px
        
        pdf_pages = []
        
        # 将图片分组，每组2张
        for page_idx in range(0, len(images), IMAGES_PER_PAGE):
            page_images = images[page_idx:page_idx + IMAGES_PER_PAGE]
            
            # 创建新页面（确保尺寸是整数）
            try:
                page = Image.new('RGB', (int(PAGE_WIDTH), int(PAGE_HEIGHT)), color='white')
            except Exception as e:
                print(f"  [错误] 创建页面失败: {e}, 尺寸: {PAGE_WIDTH}x{PAGE_HEIGHT}")
                continue
            
            # 在页面上排列图片（每页2张，上下排列）
            for idx, img in enumerate(page_images):
                row = idx // GRID_COLS
                col = idx % GRID_COLS
                
                # 计算图片在页面中的位置（使用不同的水平和垂直边距）
                x = MARGIN_H + col * (IMAGE_WIDTH + MARGIN_H)
                y = MARGIN_V + row * (IMAGE_HEIGHT + MARGIN_V)
                
                # 缩放图片以适应网格单元格（保持宽高比，使用高质量缩放算法）
                img_width, img_height = img.size
                scale_w = IMAGE_WIDTH / img_width
                scale_h = IMAGE_HEIGHT / img_height
                scale = min(scale_w, scale_h)  # 使用较小的缩放比例，确保图片完全显示
                
                new_width = int(img_width * scale)
                new_height = int(img_height * scale)
                
                # 使用高质量缩放算法（LANCZOS），保持图片清晰度，不压缩
                try:
                    # Pillow 10.0+，使用高质量重采样
                    resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                except AttributeError:
                    # Pillow < 10.0
                    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
                
                # 计算居中位置
                center_x = x + (IMAGE_WIDTH - new_width) // 2
                center_y = y + (IMAGE_HEIGHT - new_height) // 2
                
                # 将图片粘贴到页面上
                page.paste(resized_img, (center_x, center_y))
            
            pdf_pages.append(page)
            print(f"  ✓ 创建第 {len(pdf_pages)} 页（包含 {len(page_images)} 张图片）")
        
        # 保存为PDF（使用高分辨率，保持图片质量，不压缩）
        if pdf_pages:
            pdf_pages[0].save(
                output_path,
                'PDF',
                resolution=300.0,  # 提高分辨率到300 DPI，保持图片清晰度，不压缩
                save_all=True,
                append_images=pdf_pages[1:] if len(pdf_pages) > 1 else []
            )
            print(f"  ✓ 成功合并 {len(images)} 张图片到 {len(pdf_pages)} 页PDF: {output_path.name}")
            return True
        else:
            print("  ❌ 没有页面可以保存")
            return False
        
    except Exception as e:
        print(f"  ❌ 合并PDF失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def merge_images_to_pdf_img2pdf(image_paths: List[Path], output_path: Path) -> bool:
    """
    使用img2pdf将图片合并为PDF
    
    Args:
        image_paths: 图片文件路径列表
        output_path: 输出PDF文件路径
        
    Returns:
        是否成功
    """
    if not IMG2PDF_AVAILABLE:
        return False
    
    try:
        # 过滤有效的图片文件
        valid_paths = []
        for img_path in image_paths:
            if img_path.exists() and img_path.is_file():
                valid_paths.append(str(img_path))
                print(f"  ✓ 加载图片: {img_path.name}")
            else:
                print(f"  ⚠️ 跳过无效文件: {img_path.name}")
        
        if not valid_paths:
            print("  ❌ 没有有效的图片可以合并")
            return False
        
        # 使用img2pdf合并
        with open(output_path, 'wb') as f:
            f.write(img2pdf.convert(valid_paths))
        
        print(f"  ✓ 成功合并 {len(valid_paths)} 张图片到: {output_path.name}")
        return True
        
    except Exception as e:
        print(f"  ❌ 合并PDF失败: {e}")
        return False


def merge_images_to_pdf(image_paths: List[Path], output_path: Path) -> bool:
    """
    将图片合并为PDF（自动选择可用方法）
    
    Args:
        image_paths: 图片文件路径列表
        output_path: 输出PDF文件路径
        
    Returns:
        是否成功
    """
    # 优先使用Pillow，如果不可用则使用img2pdf
    if PIL_AVAILABLE:
        return merge_images_to_pdf_pillow(image_paths, output_path)
    elif IMG2PDF_AVAILABLE:
        return merge_images_to_pdf_img2pdf(image_paths, output_path)
    else:
        print("  ❌ 错误：未安装 Pillow 或 img2pdf，无法合并PDF")
        print("     请安装其中一个：pip install Pillow 或 pip install img2pdf")
        return False


def create_cover_page_with_type(employee_id: str, employee_name: str, folder_type: str) -> Image.Image | None:
    """
    创建封面页图片："工号-姓名-文件夹类型"
    
    Args:
        employee_id: 工号
        employee_name: 姓名
        folder_type: 文件夹类型（"调整前"或"调整后"）
        
    Returns:
        封面页图片对象，失败返回None
    """
    if not PIL_AVAILABLE:
        return None
    
    try:
        # PDF页面尺寸（A4横向 @ 300 DPI）
        PAGE_WIDTH = 3508
        PAGE_HEIGHT = 2480
        
        # 确保尺寸有效且为整数
        PAGE_WIDTH = int(PAGE_WIDTH)
        PAGE_HEIGHT = int(PAGE_HEIGHT)
        
        if PAGE_WIDTH <= 0 or PAGE_HEIGHT <= 0:
            print(f"  [错误] 无效的封面页尺寸: {PAGE_WIDTH}x{PAGE_HEIGHT}")
            return None
        
        # 创建白色背景图片
        try:
            img = Image.new('RGB', (PAGE_WIDTH, PAGE_HEIGHT), color='white')
        except Exception as img_error:
            print(f"  [错误] 创建封面页图片失败: {img_error}, 尺寸: {PAGE_WIDTH}x{PAGE_HEIGHT}")
            print(f"  [调试] PAGE_WIDTH类型: {type(PAGE_WIDTH)}, 值: {PAGE_WIDTH}")
            print(f"  [调试] PAGE_HEIGHT类型: {type(PAGE_HEIGHT)}, 值: {PAGE_HEIGHT}")
            import traceback
            traceback.print_exc()
            return None
        draw = ImageDraw.Draw(img)
        
        # 尝试加载字体
        font = None
        font_paths = [
            r"C:\Windows\Fonts\simhei.ttf",  # 黑体
            r"C:\Windows\Fonts\simsun.ttc",  # 宋体
            r"C:\Windows\Fonts\msyh.ttc",    # 微软雅黑
            r"C:\Windows\Fonts\msyhbd.ttc",  # 微软雅黑粗体
        ]
        
        for font_path in font_paths:
            if os.path.exists(font_path):
                try:
                    font = ImageFont.truetype(font_path, 120)  # 大号字体
                    break
                except Exception:
                    continue
        
        if font is None:
            try:
                font = ImageFont.truetype("arial.ttf", 120)
            except Exception:
                font = ImageFont.load_default()
        
        # 封面文本
        cover_text = f"{employee_id}-{employee_name}-{folder_type}"
        
        # 计算文本尺寸并居中显示
        try:
            bbox = draw.textbbox((0, 0), cover_text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
        except Exception:
            try:
                text_width, text_height = draw.textsize(cover_text, font=font)
            except Exception:
                text_width = len(cover_text) * 60
                text_height = 120
        
        x = (PAGE_WIDTH - text_width) // 2
        y = (PAGE_HEIGHT - text_height) // 2
        
        # 绘制文本
        draw.text((x, y), cover_text, fill='black', font=font)
        
        print(f"  ✓ 创建封面页: {cover_text}")
        return img
        
    except Exception as e:
        print(f"  ❌ 创建封面页失败: {e}")
        import traceback
        traceback.print_exc()
        return None


def process_folder_images_direct(image_files: Dict[str, List[Path]], folder_type: str, pdf_pages: List[Image.Image], layout_type: str = "1*2") -> None:
    """
    处理图片文件列表，添加到PDF页面列表（直接使用已获取的图片列表，避免重复扫描）
    
    Args:
        image_files: 图片文件字典，包含 'district' 和 'group' 两个列表
        folder_type: 文件夹类型（"调整前"或"调整后"）
        pdf_pages: PDF页面列表（会追加新页面）
        layout_type: 布局类型，"1*2"（每页2张，上下排列）或 "2*2"（每页4张，2x2网格）
    """
    district_images = image_files['district']  # 行政区图
    group_images = image_files['group']        # 网组网点图
    
    print(f"    [{folder_type}] 找到 {len(district_images)} 张行政区图")
    print(f"    [{folder_type}] 找到 {len(group_images)} 张网组网点图")
    
    # 1. 将行政区图单独一页
    for district_img_path in district_images:
        add_single_image_page(district_img_path, pdf_pages)
    
    # 2. 将网组网点图按指定布局合并
    if group_images:
        # 定义PDF页面尺寸（A4横向）
        PAGE_WIDTH = int(3508)
        PAGE_HEIGHT = int(2480)
        
        # 确保尺寸有效
        if PAGE_WIDTH <= 0 or PAGE_HEIGHT <= 0:
            print(f"    [错误] 无效的页面尺寸: {PAGE_WIDTH}x{PAGE_HEIGHT}")
            return
        
        # 根据布局类型设置参数
        if layout_type == "2*2":
            # 2*2布局：每页4张图片，2行2列
            IMAGES_PER_PAGE = 4
            GRID_COLS = 2
            GRID_ROWS = 2
        else:
            # 1*2布局：每页2张图片，2行1列（上下排列）
            IMAGES_PER_PAGE = 2
            GRID_COLS = 1
            GRID_ROWS = 2
        
        # 计算每张图片在页面中的尺寸
        MARGIN_H = 10
        MARGIN_V = 3
        if layout_type == "2*2":
            # 2*2布局：左右各10px，中间10px；上下各3px，中间3px
            IMAGE_WIDTH = (PAGE_WIDTH - MARGIN_H * 3) // GRID_COLS
            IMAGE_HEIGHT = (PAGE_HEIGHT - MARGIN_V * 3) // GRID_ROWS
        else:
            # 1*2布局：左右各10px；上下各3px，中间3px
            IMAGE_WIDTH = (PAGE_WIDTH - MARGIN_H * 2) // GRID_COLS
            IMAGE_HEIGHT = (PAGE_HEIGHT - MARGIN_V * 3) // GRID_ROWS
        
        # 确保计算出的尺寸有效
        if IMAGE_WIDTH <= 0 or IMAGE_HEIGHT <= 0:
            print(f"    [错误] 无效的图片尺寸: {IMAGE_WIDTH}x{IMAGE_HEIGHT}")
            return
        
        # 将图片分组，每组2张（确保每张图片只处理一次）
        for page_idx in range(0, len(group_images), IMAGES_PER_PAGE):
            # 获取当前页的图片（最多2张）
            page_image_paths = group_images[page_idx:page_idx + IMAGES_PER_PAGE]
            
            # 如果当前页没有图片，跳过
            if not page_image_paths:
                continue
            
            # 创建新页面（确保尺寸是整数）
            try:
                page = Image.new('RGB', (PAGE_WIDTH, PAGE_HEIGHT), color='white')
            except Exception as e:
                print(f"    [错误] 创建页面失败: {e}")
                print(f"    [调试] PAGE_WIDTH: {PAGE_WIDTH} (类型: {type(PAGE_WIDTH)})")
                print(f"    [调试] PAGE_HEIGHT: {PAGE_HEIGHT} (类型: {type(PAGE_HEIGHT)})")
                import traceback
                traceback.print_exc()
                continue
            
            # 在页面上排列图片（每页2张，上下排列）
            for idx, img_path in enumerate(page_image_paths):
                try:
                    # 确保索引不超过每页图片数量
                    if idx >= IMAGES_PER_PAGE:
                        break
                    
                    # 打开图片
                    try:
                        img = Image.open(img_path)
                    except Exception as open_error:
                        print(f"    [错误] 无法打开图片 {img_path.name}: {open_error}")
                        continue
                    
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # 获取图片尺寸并验证
                    img_width, img_height = img.size
                    if img_width <= 0 or img_height <= 0:
                        print(f"    [错误] 无效的图片尺寸: {img_width}x{img_height}, 文件: {img_path.name}")
                        continue
                    
                    # 计算行和列
                    if layout_type == "2*2":
                        # 2*2布局：2行2列
                        row = idx // GRID_COLS
                        col = idx % GRID_COLS
                    else:
                        # 1*2布局：2行1列（上下排列）
                        row = idx  # 因为只有1列，idx就是行号（0或1）
                        col = 0    # 只有1列，所以col始终为0
                    
                    # 计算图片在页面中的位置
                    x = MARGIN_H + col * (IMAGE_WIDTH + MARGIN_H)
                    y = MARGIN_V + row * (IMAGE_HEIGHT + MARGIN_V)
                    
                    # 缩放图片以适应网格单元格
                    scale_w = IMAGE_WIDTH / img_width
                    scale_h = IMAGE_HEIGHT / img_height
                    scale = min(scale_w, scale_h)
                    
                    new_width = int(img_width * scale)
                    new_height = int(img_height * scale)
                    
                    # 验证缩放后的尺寸
                    if new_width <= 0 or new_height <= 0:
                        print(f"    [错误] 缩放后尺寸无效: {new_width}x{new_height}, 原始: {img_width}x{img_height}, 文件: {img_path.name}")
                        continue
                    
                    # 使用高质量缩放算法
                    try:
                        resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    except AttributeError:
                        resized_img = img.resize((new_width, new_height), Image.LANCZOS)
                    
                    # 计算居中位置（确保是整数）
                    center_x = int(x + (IMAGE_WIDTH - new_width) // 2)
                    center_y = int(y + (IMAGE_HEIGHT - new_height) // 2)
                    
                    # 确保位置在页面范围内
                    center_x = max(0, min(center_x, PAGE_WIDTH - new_width))
                    center_y = max(0, min(center_y, PAGE_HEIGHT - new_height))
                    
                    # 将图片粘贴到页面上
                    page.paste(resized_img, (center_x, center_y))
                    print(f"    [OK] 加载图片 [{idx+1}/{len(page_image_paths)}]: {img_path.name}")
                except Exception as e:
                    print(f"    [警告] 跳过无效图片 {img_path.name}: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            # 只有当页面中有图片时才添加页面
            pdf_pages.append(page)
            print(f"    ✓ 创建第 {len(pdf_pages)} 页（包含 {len(page_image_paths)} 张图片）")


def add_single_image_slide_ppt(prs: Presentation, image_path: Path) -> bool:
    """
    将单张图片添加为PPT幻灯片（用于行政区图）
    
    Args:
        prs: Presentation对象
        image_path: 图片路径
        
    Returns:
        是否成功
    """
    if not PPTX_AVAILABLE:
        return False
    
    try:
        # 创建新幻灯片（空白布局）
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6是空白布局
        
        # 获取幻灯片尺寸（单位：英寸）
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 打开图片获取尺寸
        img = Image.open(image_path)
        img_width, img_height = img.size
        
        # 计算缩放比例，保持宽高比，适应幻灯片
        # 转换为英寸（假设图片DPI为96）
        img_width_inch = img_width / 96.0
        img_height_inch = img_height / 96.0
        
        slide_width_inch = slide_width / 914400.0  # PPT单位转换为英寸
        slide_height_inch = slide_height / 914400.0
        
        scale_w = slide_width_inch / img_width_inch
        scale_h = slide_height_inch / img_height_inch
        scale = min(scale_w, scale_h)  # 使用较小的缩放比例，确保图片完全显示
        
        new_width = img_width_inch * scale
        new_height = img_height_inch * scale
        
        # 计算居中位置
        left = (slide_width_inch - new_width) / 2
        top = (slide_height_inch - new_height) / 2
        
        # 添加图片到幻灯片
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(new_width), height=Inches(new_height))
        
        print(f"    ✓ 添加单页图片到PPT: {image_path.name}")
        return True
        
    except Exception as e:
        print(f"    ⚠️ 添加单页图片到PPT失败 {image_path.name}: {e}")
        return False


def process_folder_images_ppt(image_files: Dict[str, List[Path]], folder_type: str, prs: Presentation, layout_type: str = "1*2") -> None:
    """
    处理图片文件列表，添加到PPT演示文稿（直接使用已获取的图片列表）
    
    Args:
        image_files: 图片文件字典，包含 'district' 和 'group' 两个列表
        folder_type: 文件夹类型（"调整前"或"调整后"）
        prs: Presentation对象
        layout_type: 布局类型，"1*2"（每页2张，上下排列）或 "2*2"（每页4张，2x2网格）
    """
    district_images = image_files['district']  # 行政区图
    group_images = image_files['group']        # 网组网点图
    
    print(f"    [{folder_type}] 找到 {len(district_images)} 张行政区图")
    print(f"    [{folder_type}] 找到 {len(group_images)} 张网组网点图")
    
    # 1. 将行政区图单独一页
    for district_img_path in district_images:
        add_single_image_slide_ppt(prs, district_img_path)
    
    # 2. 将网组网点图按指定布局合并
    if group_images:
        # 根据布局类型设置参数
        if layout_type == "2*2":
            # 2*2布局：每页4张图片，2行2列
            IMAGES_PER_PAGE = 4
            GRID_COLS = 2
            GRID_ROWS = 2
        else:
            # 1*2布局：每页2张图片，2行1列（上下排列）
            IMAGES_PER_PAGE = 2
            GRID_COLS = 1
            GRID_ROWS = 2
        
        # 获取幻灯片尺寸（单位：英寸）
        slide_width = prs.slide_width / 914400.0  # 转换为英寸
        slide_height = prs.slide_height / 914400.0
        
        # 计算边距和图片尺寸
        MARGIN_H = 0.2  # 左右边距（英寸）
        MARGIN_V = 0.2  # 上下边距（英寸）
        
        if layout_type == "2*2":
            # 2*2布局
            IMAGE_WIDTH = (slide_width - MARGIN_H * 3) / GRID_COLS
            IMAGE_HEIGHT = (slide_height - MARGIN_V * 3) / GRID_ROWS
        else:
            # 1*2布局
            IMAGE_WIDTH = (slide_width - MARGIN_H * 2) / GRID_COLS
            IMAGE_HEIGHT = (slide_height - MARGIN_V * 3) / GRID_ROWS
        
        # 将图片分组
        for page_idx in range(0, len(group_images), IMAGES_PER_PAGE):
            page_image_paths = group_images[page_idx:page_idx + IMAGES_PER_PAGE]
            
            if not page_image_paths:
                continue
            
            # 创建新幻灯片（空白布局）
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 在幻灯片上排列图片
            for idx, img_path in enumerate(page_image_paths):
                try:
                    if idx >= IMAGES_PER_PAGE:
                        break
                    
                    # 打开图片获取尺寸
                    img = Image.open(img_path)
                    img_width, img_height = img.size
                    
                    # 转换为英寸（假设图片DPI为96）
                    img_width_inch = img_width / 96.0
                    img_height_inch = img_height / 96.0
                    
                    # 计算行和列
                    if layout_type == "2*2":
                        row = idx // GRID_COLS
                        col = idx % GRID_COLS
                    else:
                        row = idx
                        col = 0
                    
                    # 计算图片在幻灯片中的位置
                    x = MARGIN_H + col * (IMAGE_WIDTH + MARGIN_H)
                    y = MARGIN_V + row * (IMAGE_HEIGHT + MARGIN_V)
                    
                    # 计算缩放比例
                    scale_w = IMAGE_WIDTH / img_width_inch
                    scale_h = IMAGE_HEIGHT / img_height_inch
                    scale = min(scale_w, scale_h)
                    
                    new_width = img_width_inch * scale
                    new_height = img_height_inch * scale
                    
                    # 计算居中位置（在网格单元格内居中）
                    center_x = x + (IMAGE_WIDTH - new_width) / 2
                    center_y = y + (IMAGE_HEIGHT - new_height) / 2
                    
                    # 添加图片到幻灯片
                    slide.shapes.add_picture(
                        str(img_path),
                        Inches(center_x),
                        Inches(center_y),
                        width=Inches(new_width),
                        height=Inches(new_height)
                    )
                    print(f"    ✓ 加载图片到PPT [{idx+1}/{len(page_image_paths)}]: {img_path.name}")
                except Exception as e:
                    print(f"    ⚠️ 跳过无效图片 {img_path.name}: {e}")
                    continue
            
            print(f"    ✓ 创建PPT幻灯片（包含 {len(page_image_paths)} 张图片）")


def create_cover_slide_ppt(prs: Presentation, employee_id: str, employee_name: str, folder_type: str) -> bool:
    """
    创建封面幻灯片："工号-姓名-文件夹类型"
    
    Args:
        prs: Presentation对象
        employee_id: 工号
        employee_name: 姓名
        folder_type: 文件夹类型（"调整前"或"调整后"）
        
    Returns:
        是否成功
    """
    if not PPTX_AVAILABLE:
        return False
    
    try:
        # 创建新幻灯片（使用标题和内容布局）
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # 0是标题幻灯片布局
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"{employee_id}-{employee_name}-{folder_type}"
        
        # 设置标题字体大小
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        print(f"  ✓ 创建PPT封面页: {employee_id}-{employee_name}-{folder_type}")
        return True
        
    except Exception as e:
        print(f"  ❌ 创建PPT封面页失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def process_employee_folder(employee_folder: Path, pdf_pages: List[Image.Image] = None, prs: Presentation = None, layout_type: str = "1*2", output_format: str = "pdf") -> bool:
    """
    处理单个工号文件夹，将其内容添加到PDF页面列表或PPT演示文稿
    
    处理流程：
    1. 检查"调整前"和"调整后"文件夹是否存在
    2. 为每个存在的文件夹创建对应的封面页："工号-姓名-调整前"或"工号-姓名-调整后"
    3. 将每个文件夹内的行政区图图片单独一页
    4. 将其余图片（网组网点图）按指定布局合并
    
    Args:
        employee_folder: 工号-姓名文件夹路径
        pdf_pages: PDF页面列表（会追加新页面），PDF格式时使用
        prs: Presentation对象，PPT格式时使用
        layout_type: 布局类型，"1*2"（每页2张，上下排列）或 "2*2"（每页4张，2x2网格）
        output_format: 输出格式，"pdf" 或 "ppt"
        
    Returns:
        是否成功
    """
    # 解析工号和姓名
    employee_info = parse_employee_folder(employee_folder.name)
    if not employee_info:
        print(f"⚠️ 跳过无效文件夹: {employee_folder.name}（格式应为：工号-姓名）")
        return False
    
    employee_id, employee_name = employee_info
    print(f"\n处理工号: {employee_id}, 姓名: {employee_name}")
    
    # 检查"调整前"和"调整后"文件夹
    adjustment_before_folder = employee_folder / "调整前"
    adjustment_after_folder = employee_folder / "调整后"
    
    has_before = adjustment_before_folder.exists() and adjustment_before_folder.is_dir()
    has_after = adjustment_after_folder.exists() and adjustment_after_folder.is_dir()
    
    if not has_before and not has_after:
        print(f"  ⚠️ 未找到'调整前'或'调整后'文件夹，跳过")
        return False
    
    # 检查必要的库
    if output_format == "pdf":
        if not PIL_AVAILABLE:
            print("  ❌ 错误：未安装 Pillow，无法生成PDF")
            return False
    elif output_format == "ppt":
        if not PPTX_AVAILABLE:
            print("  ❌ 错误：未安装 python-pptx，无法生成PPT")
            return False
    
    try:
        # 处理"调整前"文件夹
        if has_before:
            # 获取图片文件（只获取一次）
            before_images = get_image_files_from_folder(adjustment_before_folder)
            
            # 检查是否有图片文件
            if before_images['district'] or before_images['group']:
                # 1. 创建封面页
                if output_format == "pdf":
                    cover_img = create_cover_page_with_type(employee_id, employee_name, "调整前")
                    if cover_img:
                        pdf_pages.append(cover_img)
                else:  # ppt
                    create_cover_slide_ppt(prs, employee_id, employee_name, "调整前")
                
                # 2. 处理文件夹内的图片（直接使用已获取的图片列表）
                if output_format == "pdf":
                    process_folder_images_direct(before_images, "调整前", pdf_pages, layout_type)
                else:  # ppt
                    process_folder_images_ppt(before_images, "调整前", prs, layout_type)
        
        # 处理"调整后"文件夹
        if has_after:
            # 获取图片文件（只获取一次）
            after_images = get_image_files_from_folder(adjustment_after_folder)
            
            # 检查是否有图片文件
            if after_images['district'] or after_images['group']:
                # 1. 创建封面页
                if output_format == "pdf":
                    cover_img = create_cover_page_with_type(employee_id, employee_name, "调整后")
                    if cover_img:
                        pdf_pages.append(cover_img)
                else:  # ppt
                    create_cover_slide_ppt(prs, employee_id, employee_name, "调整后")
                
                # 2. 处理文件夹内的图片（直接使用已获取的图片列表）
                if output_format == "pdf":
                    process_folder_images_direct(after_images, "调整后", pdf_pages, layout_type)
                else:  # ppt
                    process_folder_images_ppt(after_images, "调整后", prs, layout_type)
        
        print(f"  ✅ 完成: {employee_id}-{employee_name}")
        return True
        
    except Exception as e:
        print(f"  ❌ 处理失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """主函数"""
    print("=" * 60)
    print("网组网点图合并工具（所有工号合并为一个文件）")
    print("=" * 60)
    
    # 让用户选择输出格式
    print("\n请选择输出格式：")
    print("  1. PDF 格式")
    print("  2. PPT 格式")
    
    while True:
        format_choice = input("\n请输入选项 (1 或 2，默认1): ").strip()
        if not format_choice:
            format_choice = "1"
        
        if format_choice == "1":
            output_format = "pdf"
            print("✓ 已选择：PDF 格式")
            break
        elif format_choice == "2":
            if not PPTX_AVAILABLE:
                print("❌ 错误：未安装 python-pptx，无法生成PPT")
                print("   请先安装: pip install python-pptx")
                input("按回车键退出...")
                return
            output_format = "ppt"
            print("✓ 已选择：PPT 格式")
            break
        else:
            print("⚠️ 无效选项，请输入 1 或 2")
    
    print()
    
    # 让用户选择布局类型
    print("请选择合并布局：")
    print("  1. 1*2 布局（每页2张图片，上下排列）")
    print("  2. 2*2 布局（每页4张图片，2x2网格）")
    
    while True:
        layout_choice = input("\n请输入选项 (1 或 2，默认1): ").strip()
        if not layout_choice:
            layout_choice = "1"
        
        if layout_choice == "1":
            layout_type = "1*2"
            print("✓ 已选择：1*2 布局（每页2张图片，上下排列）")
            break
        elif layout_choice == "2":
            layout_type = "2*2"
            print("✓ 已选择：2*2 布局（每页4张图片，2x2网格）")
            break
        else:
            print("⚠️ 无效选项，请输入 1 或 2")
    
    print()
    
    # 获取基础目录
    base_dir = Path(get_base_dir())
    source_dir = base_dir / "网组网点路线图"
    
    # 根据输出格式设置输出目录和文件名
    if output_format == "pdf":
        output_dir = base_dir / "合并PDF"
        output_filename = "所有工号合并.pdf"
    else:  # ppt
        output_dir = base_dir / "合并PPT"
        output_filename = "所有工号合并.pptx"
    
    # 检查源目录是否存在
    if not source_dir.exists():
        print(f"❌ 错误：源目录不存在: {source_dir}")
        print("   请确保在程序目录下存在'网组网点路线图'文件夹")
        input("按回车键退出...")
        return
    
    # 创建输出目录
    output_dir.mkdir(exist_ok=True)
    print(f"输出目录: {output_dir}")
    print()
    
    # 检查必要的库
    if output_format == "pdf":
        if not PIL_AVAILABLE:
            print("❌ 错误：未安装 Pillow，无法生成PDF")
            print("   请安装：pip install Pillow")
            input("按回车键退出...")
            return
    else:  # ppt
        if not PPTX_AVAILABLE:
            print("❌ 错误：未安装 python-pptx，无法生成PPT")
            print("   请安装：pip install python-pptx")
            input("按回车键退出...")
            return
    
    # 遍历所有工号-姓名文件夹
    employee_folders = []
    for item in source_dir.iterdir():
        if item.is_dir():
            # 检查是否是工号-姓名格式的文件夹
            if parse_employee_folder(item.name):
                employee_folders.append(item)
    
    if not employee_folders:
        print("⚠️ 未找到工号-姓名格式的文件夹")
        print("   文件夹格式应为：工号-姓名，例如：FJ10331281-陈成")
        input("按回车键退出...")
        return
    
    print(f"找到 {len(employee_folders)} 个工号文件夹")
    print()
    
    # 统计信息
    total_processed = 0
    total_success = 0
    total_failed = 0
    
    # 根据输出格式初始化
    if output_format == "pdf":
        all_pdf_pages = []
    else:  # ppt
        prs = Presentation()
        # 设置幻灯片尺寸为16:9（宽屏）
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # 处理每个工号文件夹
    for employee_folder in sorted(employee_folders):
        total_processed += 1
        if output_format == "pdf":
            if process_employee_folder(employee_folder, all_pdf_pages, None, layout_type, output_format):
                total_success += 1
            else:
                total_failed += 1
        else:  # ppt
            if process_employee_folder(employee_folder, None, prs, layout_type, output_format):
                total_success += 1
            else:
                total_failed += 1
    
    # 保存文件
    output_path = output_dir / output_filename
    
    if output_format == "pdf":
        # 保存PDF文件
        if all_pdf_pages:
            try:
                all_pdf_pages[0].save(
                    output_path,
                    'PDF',
                    resolution=300.0,
                    save_all=True,
                    append_images=all_pdf_pages[1:] if len(all_pdf_pages) > 1 else []
                )
                print()
                print("=" * 60)
                print("处理完成")
                print("=" * 60)
                print(f"总计: {total_processed} 个工号")
                print(f"成功: {total_success} 个")
                print(f"失败: {total_failed} 个")
                print(f"PDF总页数: {len(all_pdf_pages)} 页")
                print(f"输出文件: {output_filename}")
                print(f"输出目录: {output_dir}")
                print()
                print("✅ 所有工号的PDF已合并为一个文件")
            except Exception as e:
                print()
                print("=" * 60)
                print("❌ 保存PDF失败")
                print("=" * 60)
                print(f"错误: {e}")
                import traceback
                traceback.print_exc()
        else:
            print()
            print("=" * 60)
            print("处理完成")
            print("=" * 60)
            print(f"总计: {total_processed} 个工号")
            print(f"成功: {total_success} 个")
            print(f"失败: {total_failed} 个")
            print()
            print("⚠️ 没有页面可以保存")
    else:  # ppt
        # 保存PPT文件
        try:
            prs.save(str(output_path))
            print()
            print("=" * 60)
            print("处理完成")
            print("=" * 60)
            print(f"总计: {total_processed} 个工号")
            print(f"成功: {total_success} 个")
            print(f"失败: {total_failed} 个")
            print(f"PPT总幻灯片数: {len(prs.slides)} 页")
            print(f"输出文件: {output_filename}")
            print(f"输出目录: {output_dir}")
            print()
            print("✅ 所有工号的PPT已合并为一个文件")
        except Exception as e:
            print()
            print("=" * 60)
            print("❌ 保存PPT失败")
            print("=" * 60)
            print(f"错误: {e}")
            import traceback
            traceback.print_exc()
    
    input("按回车键退出...")


if __name__ == "__main__":
    main()
